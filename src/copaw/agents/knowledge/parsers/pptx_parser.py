# -*- coding: utf-8 -*-
"""PPTX parser for knowledge import."""

from __future__ import annotations

from pathlib import Path

from ..exceptions import EmptyParsedContentError, KnowledgeError
from ..models import ParsedDocument


class PptxParser:
    """Parser for PPTX presentation files."""

    supported_suffixes = (".pptx",)

    def parse(self, path: Path) -> ParsedDocument:
        try:
            from pptx import (
                Presentation,
            )  # lazy import to avoid hard import cost
        except ImportError as exc:
            raise KnowledgeError(
                "python-pptx is required for PPTX knowledge import support",
            ) from exc

        presentation = Presentation(str(path))
        blocks: list[str] = []
        title = path.stem
        found_title = False
        text_block_count = 0
        table_count = 0
        notes_count = 0

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []

            for shape in slide.shapes:
                shape_lines, shape_table_count = _extract_shape_lines(shape)
                if not shape_lines:
                    continue

                if not found_title:
                    title = shape_lines[0]
                    found_title = True

                text_block_count += len(shape_lines)
                table_count += shape_table_count
                slide_lines.extend(shape_lines)

            if slide_lines:
                blocks.append(f"<<<SLIDE:{slide_index}>>>")
                blocks.extend(slide_lines)

            notes_text = _extract_notes_text(slide)
            if notes_text:
                if not found_title:
                    title = notes_text.splitlines()[0].strip()
                    found_title = True
                notes_count += 1
                blocks.append(f"<<<SLIDE:{slide_index} NOTES>>>")
                blocks.append(notes_text)

        raw_text = "\n\n".join(blocks).strip()
        if not raw_text:
            raise EmptyParsedContentError(
                "PPTX contains no extractable text content",
            )

        return ParsedDocument(
            title=title,
            source_path=str(path),
            source_type="pptx",
            raw_text=raw_text,
            metadata={
                "slide_count": len(presentation.slides),
                "text_block_count": text_block_count,
                "table_count": table_count,
                "notes_count": notes_count,
                "line_count": len(raw_text.splitlines()),
            },
        )


def _extract_shape_lines(shape) -> tuple[list[str], int]:
    if getattr(shape, "has_table", False):
        table_lines = _extract_table_lines(shape.table)
        return table_lines, 1 if table_lines else 0

    if getattr(shape, "has_text_frame", False):
        text = _normalize_text(shape.text_frame.text)
        if text:
            return [text], 0

    fallback = _normalize_text(getattr(shape, "text", ""))
    if fallback:
        return [fallback], 0
    return [], 0


def _extract_table_lines(table) -> list[str]:
    lines: list[str] = []
    for row in table.rows:
        cells = [_normalize_text(cell.text) for cell in row.cells]
        while cells and not cells[-1]:
            cells.pop()
        if not cells or not any(cells):
            continue
        lines.append(" | ".join(cells).strip())
    return lines


def _extract_notes_text(slide) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is None:
        return ""
    return _normalize_text(notes_frame.text)


def _normalize_text(value: object) -> str:
    if not value:
        return ""
    return (
        str(value)
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .replace("\v", "\n")
        .strip()
    )
